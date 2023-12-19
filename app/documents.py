import os
import PyPDF2
import base64
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores.pinecone import Pinecone
from langchain.embeddings import OpenAIEmbeddings
from io import BytesIO
from classes import DocumentType
from pptx import Presentation
from docx import Document
from bs4 import BeautifulSoup
import eventlet
eventlet.monkey_patch()


async def insert_documents_to_pinecone(documents):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=0
    )

    for document in documents:
        document_bytes = base64.b64decode(document["Bytes"])
        if document["Type"] == DocumentType.PDF.value:
            document_text = __extract_text_from_pdf_bytes(document_bytes)
        elif document["Type"] == DocumentType.PPTX.value:
            document_text = __extract_text_from_pptx_bytes(document_bytes)
        elif document["Type"] == DocumentType.DOCX.value:
            document_text = __extract_text_from_docx_bytes(document_bytes)
        elif document["Type"] == DocumentType.HTML.value:
            document_text = __extract_text_from_html_bytes(document_bytes)
        elif document["Type"] == DocumentType.TXT.value:
            document_text = __extract_text_from_txt_bytes(document_bytes)
        else:
            continue

        chunks = text_splitter.create_documents([document_text])

        try:
            index = Pinecone.get_pinecone_index(
                index_name=os.environ["INDEX_NAME"])
            vector_store = Pinecone(
                index=index, embedding=OpenAIEmbeddings(), text_key="text")

            vectors = await vector_store.aadd_documents(documents=chunks)

            document["VectorIds"] = vectors
        except Exception as e:
            print(f"An error occurred: {e}")

    return documents


def remove_documents_from_pinecone(documents):
    for document in documents:
        try:
            index = Pinecone.get_pinecone_index(
                index_name=os.environ["INDEX_NAME"])
            vector_store = Pinecone(
                index=index, embedding=OpenAIEmbeddings(), text_key="text")

            vector_store.delete(document["VectorIds"])

        except Exception as e:
            print(f"An error occurred: {e}")


def __extract_text_from_pptx_bytes(pptx_bytes):
    prs = Presentation(BytesIO(pptx_bytes))
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text


def __extract_text_from_docx_bytes(docx_bytes):
    doc = Document(BytesIO(docx_bytes))
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text


def __extract_text_from_html_bytes(html_bytes):
    soup = BeautifulSoup(html_bytes, 'html.parser')
    return soup.get_text()


def __extract_text_from_txt_bytes(txt_bytes):
    return txt_bytes.decode('utf-8')


def __extract_text_from_pdf_bytes(pdf_bytes):
    text = ''
    with BytesIO(pdf_bytes) as file:
        pdf_reader = PyPDF2.PdfReader(file)
        num_pages = len(pdf_reader.pages)
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
    return text
