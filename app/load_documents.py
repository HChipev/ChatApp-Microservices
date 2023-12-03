from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Pinecone
from langchain.embeddings import OpenAIEmbeddings
from io import BytesIO
from classes import FileType
from pptx import Presentation
from docx import Document
from bs4 import BeautifulSoup
import fitz

def insert_documents_to_pinecone(index_name, documents):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=0
    )

    for document in documents:
        
        if document.doc_type == FileType.PDF:
            document_text = __extract_text_from_pdf_bytes(document.content)
        elif document.doc_type == FileType.PPTX:
            document_text = __extract_text_from_pptx_bytes(document.content)
        elif document.doc_type == FileType.DOCX:
            document_text = __extract_text_from_docx_bytes(document.content)
        elif document.doc_type == FileType.HTML:
            document_text = __extract_text_from_html_bytes(document.content)
        elif document.doc_type == FileType.TXT:
            document_text = __extract_text_from_txt_bytes(document.content)
        else:
            continue

        chunks = text_splitter.create_documents([document_text])

        Pinecone.from_documents(chunks,OpenAIEmbeddings(),index_name=index_name)
    

def __extract_text_from_pptx_bytes(pptx_bytes):
    prs = Presentation(BytesIO(pptx_bytes))
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def __extract_text_from_docx_bytes(docx_bytes):
    doc = Document(BytesIO(docx_bytes))
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text

def __extract_text_from_html_bytes(html_bytes):
    soup = BeautifulSoup(html_bytes, 'html.parser')
    return soup.get_text()

def __extract_text_from_txt_bytes(txt_bytes):
    return txt_bytes.decode('utf-8')

def __extract_text_from_pdf_bytes(pdf_bytes):
    pdf_document = fitz.open("pdf", pdf_bytes)
    text = ''
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        text += page.get_text()
    return text